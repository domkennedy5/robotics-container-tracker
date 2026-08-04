"""
wbr_pptx.py — WBR PowerPoint slide generator.

Mirrors the PDF layout (wbr_pdf.py forensic spec):
  - Widescreen landscape (13.33" × 7.5")
  - Dark header + orange separator + title
  - 6 KPI tiles (current week) with SLA pass/fail color coding
  - 6 trend charts (python-pptx native line charts), 2 rows × 3
  - 10-row performance table with SLA inset borders on current week
  - Footer with date + time
"""
from __future__ import annotations
import io
from datetime import date
from typing import Optional, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors (matching wbr_pdf.py) ──────────────────────────────────────────────
RGB_DARK    = RGBColor(0x33, 0x33, 0x33)
RGB_ORANGE  = RGBColor(0xE8, 0xA8, 0x38)
RGB_TEAL    = RGBColor(0x4B, 0xAC, 0xC6)
RGB_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RGB_BLACK   = RGBColor(0x00, 0x00, 0x00)
RGB_GREEN   = RGBColor(0x00, 0xAF, 0x4F)
RGB_RED     = RGBColor(0xCC, 0x00, 0x00)
RGB_YELLOW  = RGBColor(0xFF, 0xF2, 0xCC)
RGB_GRAY    = RGBColor(0x66, 0x66, 0x66)
RGB_DARK_BG = RGBColor(0x1E, 0x3A, 0x5F)   # tile bg for non-SLA metrics
RGB_TILE_LBL = RGBColor(0xBB, 0xBB, 0xBB)

SLA_TARGET = 95

ROW_LABELS = [
    "Containers",
    "AV\u2192OA Avg (Days)", "AV\u2192OA SLA %", "AV\u2192OA P90 (Days)",
    "OA\u2192Del Avg (Days)", "OA\u2192Del SLA %", "OA\u2192Del P90 (Days)",
    "Empty\u2192Term SLA %", "E2E Transit Avg (Days)", "On-Time to Promise %",
]
ROW_KEYS = [
    "containers", "av_oa_avg", "av_oa_sla_pct", "av_oa_p90",
    "oa_del_avg", "oa_del_sla_pct", "oa_del_p90",
    "empty_term_pct", "e2e_avg", "otp_pct",
]
SLA_ROWS = {2, 5, 7, 9}

CHART_CONFIGS = [
    # (title, key, is_pct, row, band)
    ("AV to OA (avg. days)",       "av_oa_avg",      False, 1, 0),
    ("OA to Delivery (avg. days)", "oa_del_avg",      False, 1, 1),
    ("E2E Transit (avg. days)",    "e2e_avg",         False, 1, 2),
    ("Empty to Termination (%)",   "empty_term_pct",  True,  2, 0),
    ("On-Time to Promise %",       "otp_pct",         True,  2, 1),
    ("Volume (containers)",        "containers",      False, 2, 2),
]


def _fmt(v, key: str) -> str:
    if v is None:
        return "\u2013"
    if key in ("av_oa_sla_pct", "oa_del_sla_pct", "empty_term_pct", "otp_pct"):
        return f"{int(round(float(v)))}%"
    if key in ("av_oa_avg", "oa_del_avg", "e2e_avg"):
        return f"{float(v):.1f}"
    return str(int(round(float(v))))


def _solid_bg(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def _set_cell(cell, text, bold=False, font_size=9, font_color=RGB_BLACK,
              bg=None, align=PP_ALIGN.CENTER):
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _cell_border_rgb(tc_elem, rgb_hex: str):
    """Apply a colored inset border to a table cell via raw XML (SLA indicator)."""
    r, g, b = int(rgb_hex[0:2], 16), int(rgb_hex[2:4], 16), int(rgb_hex[4:6], 16)
    hex_str = f"{r:02X}{g:02X}{b:02X}"
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tcPr = tc_elem.find(qn("a:tcPr"))
        if tcPr is None:
            continue
        ln_elem = tcPr.find(qn(edge))
        if ln_elem is None:
            ln_elem = etree.SubElement(tcPr, qn(edge))
        ln_elem.set("w", "25400")   # ~2pt in EMU/100
        solidFill = ln_elem.find(qn("a:solidFill"))
        if solidFill is None:
            solidFill = etree.SubElement(ln_elem, qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", hex_str)


def _add_line_chart(slide, title: str, values: list, week_labels: list,
                    x, y, cx, cy, is_pct: bool = False):
    """Add a native python-pptx line chart."""
    chart_data = ChartData()
    chart_data.categories = week_labels
    # Fill None values with 0 for chart (avoids gaps)
    chart_vals = tuple(v if v is not None else None for v in values)
    chart_data.add_series(title, chart_vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart

    # Style the chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True

    # Format the series line (teal)
    s = chart.series[0]
    s.format.line.color.rgb = RGB_TEAL
    s.format.line.width = Pt(2)

    # Value axis: start at 0, hide minor gridlines
    try:
        va = chart.value_axis
        va.minimum_scale = 0
        if is_pct:
            va.maximum_scale = 100
    except Exception:
        pass

    return chart_shape


def generate_wbr_pptx(
    week_labels: List[str],
    weeks_data: List[Optional[dict]],
    totals: Optional[dict],
    report_date: date,
    report_time: str = "10:00am (CT) | 9:00am (MT) | 8:00am (PT)",
    sites_str: str = "USORF, USBOS, USSAV, USLAX",
) -> bytes:
    """Generate a WBR PowerPoint slide (mirrors PDF layout).

    Returns raw .pptx bytes for st.download_button.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    W = prs.slide_width
    H = prs.slide_height
    year_str  = str(report_date.year)
    year_2d   = year_str[2:]
    curr      = weeks_data[-1] if weeks_data and weeks_data[-1] else {}

    # ── Header bar ────────────────────────────────────────────────────────────
    HDR_H = Inches(0.7)
    hdr   = slide.shapes.add_shape(1, 0, 0, W, HDR_H)
    _solid_bg(hdr, RGB_DARK); _no_line(hdr)

    title_box = slide.shapes.add_textbox(Inches(0.25), Pt(6), Inches(9.5), Inches(0.55))
    tf = title_box.text_frame
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = "Robotics Destination Dray Metrics \u2014 NA"
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RGB_WHITE

    wk_box = slide.shapes.add_textbox(Inches(9.8), Pt(12), Inches(3.3), Inches(0.4))
    tf2 = wk_box.text_frame
    p2  = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2  = p2.add_run()
    r2.text = f"{week_labels[-1] if week_labels else ''} · {report_date.strftime('%b %d, %Y')}"
    r2.font.size = Pt(9); r2.font.color.rgb = RGB_ORANGE

    # ── Orange separator ─────────────────────────────────────────────────────
    sep = slide.shapes.add_shape(1, 0, HDR_H, W, Pt(3))
    _solid_bg(sep, RGB_ORANGE); _no_line(sep)

    # ── KPI tiles (current week) ──────────────────────────────────────────────
    TILE_TOP = HDR_H + Pt(5)
    TILE_H   = Inches(0.85)
    TILE_W   = Inches(2.1)
    TILE_GAP = Inches(0.06)
    TILE_X0  = Inches(0.15)

    kpi_defs = [
        ("CONTAINERS",    "containers",     False),
        ("AV\u2192OA SLA", "av_oa_sla_pct",  True),
        ("AV\u2192OA AVG", "av_oa_avg",      False),
        ("OA\u2192DEL SLA","oa_del_sla_pct", True),
        ("E2E AVG (d)",    "e2e_avg",        False),
        ("OTP %",          "otp_pct",        True),
    ]

    for i, (label, key, is_sla) in enumerate(kpi_defs):
        v = curr.get(key)
        val_str = _fmt(v, key)
        tx = TILE_X0 + i * (TILE_W + TILE_GAP)

        if is_sla and v is not None:
            tile_bg  = RGBColor(0x00, 0x7A, 0x1A) if float(v) >= SLA_TARGET else RGBColor(0xAA, 0x00, 0x00)
            val_clr  = RGB_WHITE
        else:
            tile_bg  = RGB_DARK_BG
            val_clr  = RGB_TEAL

        rect = slide.shapes.add_shape(1, tx, TILE_TOP, TILE_W, TILE_H)
        _solid_bg(rect, tile_bg)
        rect.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
        rect.line.width = Pt(0.5)

        lbl_box = slide.shapes.add_textbox(tx + Pt(3), TILE_TOP + Pt(3), TILE_W - Pt(6), Pt(13))
        _p = lbl_box.text_frame.paragraphs[0]
        _p.alignment = PP_ALIGN.CENTER
        _r = _p.add_run(); _r.text = label
        _r.font.size = Pt(7); _r.font.color.rgb = RGB_TILE_LBL

        val_box = slide.shapes.add_textbox(tx + Pt(2), TILE_TOP + Pt(17), TILE_W - Pt(4), Pt(30))
        _vp = val_box.text_frame.paragraphs[0]
        _vp.alignment = PP_ALIGN.CENTER
        _vr = _vp.add_run(); _vr.text = val_str
        _vr.font.size = Pt(19); _vr.font.bold = True; _vr.font.color.rgb = val_clr

    # ── Charts — 2 rows × 3 ──────────────────────────────────────────────────
    CH_TOP_R1 = TILE_TOP + TILE_H + Pt(6)
    CH_H      = Inches(1.45)
    CH_W      = Inches(4.2)
    CH_GAP    = Inches(0.1)
    CH_X0     = Inches(0.15)
    CH_TOP_R2 = CH_TOP_R1 + CH_H + Pt(6)

    chart_configs = [
        ("AV to OA (avg days)",       "av_oa_avg",      False, CH_TOP_R1),
        ("OA to Delivery (avg days)", "oa_del_avg",      False, CH_TOP_R1),
        ("E2E Transit (avg days)",    "e2e_avg",         False, CH_TOP_R1),
        ("Empty to Term %",           "empty_term_pct",  True,  CH_TOP_R2),
        ("On-Time to Promise %",      "otp_pct",         True,  CH_TOP_R2),
        ("Volume (containers)",       "containers",      False, CH_TOP_R2),
    ]

    for ci, (title, key, is_pct, row_top) in enumerate(chart_configs):
        col = ci % 3
        cx  = CH_X0 + col * (CH_W + CH_GAP)
        vals = [d.get(key) if d else None for d in weeks_data]
        _add_line_chart(slide, title, vals, week_labels, cx, row_top, CH_W, CH_H, is_pct)

    # ── Performance table ─────────────────────────────────────────────────────
    TBL_TOP  = CH_TOP_R2 + CH_H + Pt(8)
    n_weeks  = len(week_labels)
    n_cols   = 1 + n_weeks + 1    # label + weeks + total
    n_rows   = 1 + len(ROW_KEYS)  # header + data
    LABEL_W  = Inches(2.0)
    WEEK_W   = Inches(1.45)
    TOTAL_W  = Inches(1.15)
    ROW_H    = Pt(16)
    TBL_W    = LABEL_W + n_weeks * WEEK_W + TOTAL_W
    TBL_H    = n_rows * ROW_H

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.15), TBL_TOP, TBL_W, TBL_H)
    tbl = tbl_shape.table

    # Column widths
    tbl.columns[0].width = LABEL_W
    for ci in range(1, n_weeks + 1):
        tbl.columns[ci].width = WEEK_W
    tbl.columns[n_cols - 1].width = TOTAL_W

    # Row heights
    for ri in range(n_rows):
        tbl.rows[ri].height = ROW_H

    # Header row
    _set_cell(tbl.cell(0, 0), "Metric", bold=True, font_size=8, align=PP_ALIGN.LEFT)
    for ci, wk in enumerate(week_labels):
        wk_disp = f"{year_2d} {wk}"
        _set_cell(tbl.cell(0, ci + 1), wk_disp, bold=True, font_size=7)
    _set_cell(tbl.cell(0, n_cols - 1), "Total", bold=True, font_size=8, bg=RGB_YELLOW)

    # Data rows
    last_wi = n_weeks - 1
    for ri, (label, key) in enumerate(zip(ROW_LABELS, ROW_KEYS)):
        _set_cell(tbl.cell(ri + 1, 0), label, bold=True, font_size=7, align=PP_ALIGN.LEFT)
        for ci, wk_data in enumerate(weeks_data):
            v = wk_data.get(key) if wk_data else None
            txt = _fmt(v, key)
            is_last = (ci == last_wi)
            bg_c = None
            cell = tbl.cell(ri + 1, ci + 1)
            _set_cell(cell, txt, bold=is_last, font_size=7)
            # SLA border on current week
            if ri in SLA_ROWS and is_last and v is not None:
                sla_pass = float(v) >= SLA_TARGET
                clr_hex  = "00AF4F" if sla_pass else "CC0000"
                try:
                    _cell_border_rgb(cell._tc, clr_hex)
                except Exception:
                    pass
        # Total col
        tv = totals.get(key) if totals else None
        _set_cell(tbl.cell(ri + 1, n_cols - 1), _fmt(tv, key), font_size=7, bg=RGB_YELLOW)

    # ── Footer ────────────────────────────────────────────────────────────────
    try:
        date_str = report_date.strftime("%-B %-d, %Y")
    except ValueError:
        date_str = report_date.strftime("%B %d, %Y").replace(" 0", " ")

    ft = slide.shapes.add_textbox(Inches(0.2), H - Pt(18), W - Inches(0.4), Pt(16))
    ftp = ft.text_frame.paragraphs[0]
    ftr = ftp.add_run()
    ftr.text = f"{date_str}  {report_time}   \u25a0 Powered by Amazon Quick"
    ftr.font.size = Pt(7); ftr.font.color.rgb = RGB_GRAY

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
